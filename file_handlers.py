import os
from typing import Tuple, List
import PyPDF2
from docx import Document
from pptx import Presentation
import json
from bs4 import BeautifulSoup
from langchain.text_splitter import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import io

class FileHandler:
    def __init__(self):
        load_dotenv()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=int(os.getenv("CHUNK_SIZE", "4000")),
            chunk_overlap=int(os.getenv("CHUNK_OVERLAP", "200")),
            length_function=len,
        )
        self.MAX_FILE_SIZE = int(os.getenv("MAX_FILE_SIZE", "26214400"))  # Default 25MB
        self.supported_extensions = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx,
            '.json': self._extract_json,
            '.html': self._extract_html,
            '.txt': self._extract_txt,
            '.md': self._extract_txt
        }

    def validate_file(self, file_path: str) -> bool:
        """Validate file size and type."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_size = os.path.getsize(file_path)
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds the maximum limit of 25MB. Current size: {file_size / (1024*1024):.2f}MB")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        return True

    def process_file(self, file_path: str) -> List[str]:
        """Process the uploaded file and return extracted text chunks."""
        # Validate file before processing
        self.validate_file(file_path)
        
        file_ext = os.path.splitext(file_path)[1].lower()
        # Extract text using appropriate handler
        text = self.supported_extensions[file_ext](file_path)
        
        # Split text into chunks if needed
        chunks = self.text_splitter.split_text(text)
        return chunks

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text and image info from PPTX file."""
        prs = Presentation(file_path)
        self._pptx_images = []  # Store image info for preservation
        text = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
                # Preserve images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    self._pptx_images.append({
                        "slide_idx": slide_idx,
                        "shape_idx": shape_idx,
                        "image_bytes": image_bytes,
                        "image_ext": image_ext,
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height
                    })
        return "\n".join(text)

    def _extract_json(self, file_path: str) -> str:
        """Extract text from JSON file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        return json.dumps(data, ensure_ascii=False, indent=2)

    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
        return soup.get_text()

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def save_translated_file(self, translated_text: str, original_file: str, translation_details: dict = None) -> str:
        """
        Save the translated text to a file with the same format as the original.
        If translation_details is provided, only the final translation (from 'Terminology Check') for each chunk is saved.
        """
        # Get the file extension
        _, ext = os.path.splitext(original_file)
        
        # Create output filename
        output_file = original_file.replace(ext, f"_translated{ext}")
        
        # Prepare the text to save
        if translation_details and 'chunks' in translation_details:
            # Extract only the final translation for each chunk
            paragraphs = []
            for chunk in translation_details['chunks']:
                for step in chunk['steps']:
                    if step['step'] == 'Terminology Check':
                        paragraphs.append(step['result'])
                        break
            text_to_save = paragraphs
        else:
            # Fallback: use the provided translated_text as a single block
            text_to_save = [translated_text]
        
        # Save based on file type
        if ext.lower() == '.docx':
            doc = Document()
            for para in text_to_save:
                doc.add_paragraph(para)
            doc.save(output_file)
        elif ext.lower() == '.json':
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump({"translation": '\n\n'.join(text_to_save)}, f, ensure_ascii=False, indent=2)
        elif ext.lower() == '.html':
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"<html><body>{'<br><br>'.join(text_to_save)}</body></html>")
        elif ext.lower() == '.pptx':
            prs = Presentation(original_file)
            # Remove all text and images from slides
            for slide in prs.slides:
                for shape in list(slide.shapes):
                    sp_type = getattr(shape, 'shape_type', None)
                    if hasattr(shape, 'text') or sp_type == MSO_SHAPE_TYPE.PICTURE:
                        slide.shapes._spTree.remove(shape._element)
            # Add translated text and restore images
            for i, para in enumerate(text_to_save):
                if i < len(prs.slides):
                    slide = prs.slides[i]
                    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5)).text = para
            # Restore images
            if hasattr(self, '_pptx_images'):
                for img in self._pptx_images:
                    slide = prs.slides[img["slide_idx"]]
                    image_stream = io.BytesIO(img["image_bytes"])
                    slide.shapes.add_picture(
                        image_stream,
                        img["left"],
                        img["top"],
                        img["width"],
                        img["height"]
                    )
            prs.save(output_file)
        else:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(text_to_save))
        
        return output_file

    def save_translation_details(self, details: dict, original_file: str) -> str:
        """
        Save the translation details (including original text, reflections, and final translation)
        to a JSON file.
        
        Args:
            details: Dictionary containing translation details
            original_file: Path to the original file
            
        Returns:
            Path to the saved details file
        """
        # Create output filename
        output_file = original_file.replace(os.path.splitext(original_file)[1], "_details.json")
        
        # Save as JSON
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(details, f, ensure_ascii=False, indent=2)
        
        return output_file 