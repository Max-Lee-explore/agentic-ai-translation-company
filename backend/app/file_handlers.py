import os
from typing import Tuple, List
import PyPDF2
from docx import Document
from pptx import Presentation
import json
from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv

class FileHandler:
    def __init__(self):
        load_dotenv()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=int(os.getenv("CHUNK_SIZE", "4000")),
            chunk_overlap=int(os.getenv("CHUNK_OVERLAP", "200")),
            length_function=len,
        )
        self.MAX_FILE_SIZE = int(os.getenv("MAX_FILE_SIZE", "26214400"))  # Default 25MB
        self.supported_extensions = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx,
            '.json': self._extract_json,
            '.html': self._extract_html,
            '.txt': self._extract_txt,
            '.md': self._extract_txt
        }

    def validate_file(self, file_path: str) -> bool:
        """Validate file size and type."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_size = os.path.getsize(file_path)
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds the maximum limit of 25MB. Current size: {file_size / (1024*1024):.2f}MB")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        return True

    def process_file(self, file_path: str) -> List[str]:
        """Process the uploaded file and return extracted text chunks."""
        # Validate file before processing
        self.validate_file(file_path)
        
        file_ext = os.path.splitext(file_path)[1].lower()
        # Extract text using appropriate handler
        text = self.supported_extensions[file_ext](file_path)
        
        # Split text into chunks if needed
        chunks = self.text_splitter.split_text(text)
        return chunks

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_json(self, file_path: str) -> str:
        """Extract text from JSON file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        return json.dumps(data, ensure_ascii=False, indent=2)

    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
        return soup.get_text()

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text file."""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def save_translated_file(self, translated_text: str, original_file: str, translation_details: dict = None, output_format: str = None) -> str:
        """
        Save the translated text to a file with the same format as the original.
        If translation_details is provided, the final step result for each chunk is saved.
        """
        # Get the file extension
        _, original_ext = os.path.splitext(original_file)
        
        # Determine output extension
        if output_format:
            ext = f".{output_format.lower()}" if not output_format.startswith('.') else output_format.lower()
        else:
            ext = original_ext
        
        # Create output filename
        output_file = original_file.replace(original_ext, f"_translated{ext}")
        
        # Prepare the text to save
        if translation_details and 'chunks' in translation_details:
            # Extract the final translation for each chunk
            paragraphs = []
            for chunk in translation_details['chunks']:
                # Try to get Terminology Check first, otherwise use the last step
                final_result = None
                for step in chunk['steps']:
                    if step['step'] == 'Terminology Check':
                        final_result = step['result']
                        break
                
                # If no Terminology Check, use the last step (Improvement)
                if not final_result and chunk['steps']:
                    final_result = chunk['steps'][-1]['result']
                
                if final_result:
                    paragraphs.append(final_result)
            
            text_to_save = paragraphs if paragraphs else [translated_text]
        else:
            # Fallback: use the provided translated_text as a single block
            text_to_save = [translated_text]
        
        # Save based on file type
        if ext.lower() == '.docx':
            doc = Document()
            for para in text_to_save:
                doc.add_paragraph(para)
            doc.save(output_file)
        elif ext.lower() == '.json':
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump({"translation": '\n\n'.join(text_to_save)}, f, ensure_ascii=False, indent=2)
        elif ext.lower() == '.html':
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"<html><body>{'<br><br>'.join(text_to_save)}</body></html>")
        else:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(text_to_save))
        
        return output_file

    def save_translation_details(self, details: dict, original_file: str) -> str:
        """
        Save the translation details (including original text, reflections, and final translation)
        to a JSON file.
        
        Args:
            details: Dictionary containing translation details
            original_file: Path to the original file
            
        Returns:
            Path to the saved details file
        """
        # Create output filename
        output_file = original_file.replace(os.path.splitext(original_file)[1], "_details.json")
        
        # Save as JSON
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(details, f, ensure_ascii=False, indent=2)
        
        return output_file

    def parse_glossary(self, file_path: str) -> dict:
        """
        Parse a glossary file (Excel or CSV) into a dictionary.
        Expected format: First row contains language names, subsequent rows contain terms.
        Returns a dictionary where keys are source terms and values are target terms.
        """
        import pandas as pd
        
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path)
            elif ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                raise ValueError("Unsupported glossary format. Please use .xlsx, .xls, or .csv")
            
            # Basic validation: need at least 2 columns
            if len(df.columns) < 2:
                raise ValueError("Glossary must have at least 2 columns (Source Language and Target Language)")
            
            # Convert to dictionary (assuming first column is source, second is target for now)
            # In a real scenario, we might want to map specific language names
            glossary = dict(zip(df.iloc[:, 0], df.iloc[:, 1]))
            
            # Clean up dictionary (remove NaNs, empty strings)
            return {str(k).strip(): str(v).strip() for k, v in glossary.items() if pd.notna(k) and pd.notna(v)}
            
        except Exception as e:
            print(f"Error parsing glossary: {e}")
            raise 